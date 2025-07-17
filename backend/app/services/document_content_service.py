"""Document content processing service."""

import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


class DocumentContentService:
    """文档内容处理服务."""

    def __init__(self) -> None:
        """初始化文档服务."""
        self.supported_types = {
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".doc": self._extract_doc,
            ".txt": self._extract_txt,
            ".md": self._extract_markdown,
            ".pptx": self._extract_pptx,
            ".ppt": self._extract_ppt,
        }

        # 尝试初始化markitdown（可选增强）
        try:
            from markitdown import MarkItDown
            self.markitdown: Any = MarkItDown()
            logger.info("MarkItDown initialized successfully for enhanced document conversion")
        except ImportError:
            self.markitdown = None
            logger.info("MarkItDown not available, using default extractors")

    async def extract_content(self, file_path: Path, file_ext: str) -> str:
        """提取文档内容."""
        try:
            # 如果markitdown可用且支持该格式，优先使用
            if self.markitdown and file_ext in ['.pdf', '.docx', '.pptx', '.xlsx']:
                try:
                    result = self.markitdown.convert(str(file_path))
                    if result and result.text_content:
                        logger.info(f"Successfully extracted {file_path} using MarkItDown")
                        return result.text_content
                except Exception as e:
                    logger.warning(f"MarkItDown extraction failed for {file_path}: {e}, falling back to default")

            # 使用默认提取器
            if file_ext not in self.supported_types:
                raise ValueError(f"不支持的文件类型: {file_ext}")

            extractor = self.supported_types[file_ext]
            content = await extractor(file_path)

            return content or ""

        except Exception as e:
            logger.error(f"提取文档内容失败 {file_path}: {str(e)}")
            raise

    async def _extract_pdf(self, file_path: Path) -> str:
        """提取PDF内容."""
        try:
            import pdfplumber

            content = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        content.append(text)

                    # 提取表格
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            # 将表格转换为文本
                            for row in table:
                                content.append(" | ".join(str(cell) if cell else "" for cell in row))

            return "\n".join(content)

        except ImportError:
            logger.warning("pdfplumber 未安装，使用基础文本提取")
            return await self._extract_with_textract(file_path)
        except Exception as e:
            logger.error(f"PDF提取失败: {str(e)}")
            return await self._extract_with_textract(file_path)

    async def _extract_docx(self, file_path: Path) -> str:
        """提取DOCX内容."""
        try:
            from docx import Document

            doc = Document(str(file_path))
            content = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)

            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content.append(" | ".join(row_text))

            return "\n".join(content)

        except ImportError:
            logger.warning("python-docx 未安装，使用基础文本提取")
            return await self._extract_with_textract(file_path)
        except Exception as e:
            logger.error(f"DOCX提取失败: {str(e)}")
            return await self._extract_with_textract(file_path)

    async def _extract_doc(self, file_path: Path) -> str:
        """提取DOC内容."""
        try:
            import textract  # type: ignore[import-untyped]

            # 使用textract处理老版本Word文档
            text = textract.process(str(file_path))
            return text.decode("utf-8")

        except ImportError:
            logger.warning("textract 未安装，无法处理DOC文件")
            return "无法提取DOC文件内容，请转换为DOCX格式"
        except Exception as e:
            logger.error(f"DOC提取失败: {str(e)}")
            return "DOC文件提取失败"

    async def _extract_txt(self, file_path: Path) -> str:
        """提取TXT内容."""
        try:
            # 尝试多种编码
            encodings = ["utf-8", "gbk", "gb2312", "latin-1"]

            for encoding in encodings:
                try:
                    with open(file_path, encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue

            # 如果所有编码都失败，使用错误处理
            with open(file_path, encoding="utf-8", errors="ignore") as file:
                return file.read()

        except Exception as e:
            logger.error(f"TXT提取失败: {str(e)}")
            return ""

    async def _extract_markdown(self, file_path: Path) -> str:
        """提取Markdown内容."""
        return await self._extract_txt(file_path)

    async def _extract_pptx(self, file_path: Path) -> str:
        """提取PPTX内容."""
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            content = []

            for slide in prs.slides:
                slide_content = []

                for shape in slide.shapes:
                    try:
                        if shape.text.strip():  # type: ignore[attr-defined]
                            slide_content.append(shape.text)  # type: ignore[attr-defined]
                    except AttributeError:
                        continue

                if slide_content:
                    content.append("\n".join(slide_content))

            return "\n\n---\n\n".join(content)

        except ImportError:
            logger.warning("python-pptx 未安装，使用基础文本提取")
            return await self._extract_with_textract(file_path)
        except Exception as e:
            logger.error(f"PPTX提取失败: {str(e)}")
            return await self._extract_with_textract(file_path)

    async def _extract_ppt(self, file_path: Path) -> str:
        """提取PPT内容."""
        try:
            import textract  # type: ignore[import-untyped]

            text = textract.process(str(file_path))
            return text.decode("utf-8")

        except ImportError:
            logger.warning("textract 未安装，无法处理PPT文件")
            return "无法提取PPT文件内容，请转换为PPTX格式"
        except Exception as e:
            logger.error(f"PPT提取失败: {str(e)}")
            return "PPT文件提取失败"

    async def _extract_with_textract(self, file_path: Path) -> str:
        """使用textract提取内容."""
        try:
            import textract  # type: ignore[import-untyped]

            text = textract.process(str(file_path))
            return text.decode("utf-8")

        except ImportError:
            logger.warning("textract 未安装，无法提取文档内容")
            return "无法提取文档内容，请安装相应的处理库"
        except Exception as e:
            logger.error(f"textract提取失败: {str(e)}")
            return "文档内容提取失败"

    async def get_document_metadata(self, file_path: Path) -> dict[str, Any]:
        """获取文档元数据."""
        try:
            stat = file_path.stat()

            metadata: dict[str, Any] = {
                "file_size": stat.st_size,
                "created_time": stat.st_ctime,
                "modified_time": stat.st_mtime,
                "file_extension": file_path.suffix.lower(),
            }

            # 根据文件类型获取特定元数据
            file_ext = file_path.suffix.lower()

            if file_ext == ".pdf":
                metadata.update(await self._get_pdf_metadata(file_path))
            elif file_ext in [".docx", ".doc"]:
                metadata.update(await self._get_doc_metadata(file_path))
            elif file_ext in [".pptx", ".ppt"]:
                metadata.update(await self._get_ppt_metadata(file_path))

            return metadata

        except Exception as e:
            logger.error(f"获取文档元数据失败: {str(e)}")
            return {}

    async def _get_pdf_metadata(self, file_path: Path) -> dict[str, Any]:
        """获取PDF元数据."""
        try:
            import pdfplumber

            with pdfplumber.open(file_path) as pdf:
                metadata: dict[str, Any] = {
                    "page_count": len(pdf.pages),
                }

                # pdfplumber的元数据访问方式
                if pdf.metadata:
                    metadata.update(
                        {
                            "title": pdf.metadata.get("Title", ""),
                            "author": pdf.metadata.get("Author", ""),
                            "subject": pdf.metadata.get("Subject", ""),
                            "creator": pdf.metadata.get("Creator", ""),
                            "producer": pdf.metadata.get("Producer", ""),
                            "creation_date": str(pdf.metadata.get("CreationDate", "")),
                            "modification_date": str(pdf.metadata.get("ModDate", "")),
                        }
                    )

                return metadata

        except Exception as e:
            logger.error(f"获取PDF元数据失败: {str(e)}")
            return {}

    async def _get_doc_metadata(self, file_path: Path) -> dict[str, Any]:
        """获取Word文档元数据."""
        try:
            from docx import Document

            if file_path.suffix.lower() == ".docx":
                doc = Document(str(file_path))

                metadata: dict[str, Any] = {
                    "paragraph_count": len(doc.paragraphs),
                    "table_count": len(doc.tables),
                }

                if doc.core_properties:
                    metadata.update(
                        {
                            "title": doc.core_properties.title or "",
                            "author": doc.core_properties.author or "",
                            "subject": doc.core_properties.subject or "",
                            "created": str(doc.core_properties.created)
                            if doc.core_properties.created
                            else "",
                            "modified": str(doc.core_properties.modified)
                            if doc.core_properties.modified
                            else "",
                        }
                    )

                return metadata

            return {}

        except Exception as e:
            logger.error(f"获取Word元数据失败: {str(e)}")
            return {}

    async def _get_ppt_metadata(self, file_path: Path) -> dict[str, Any]:
        """获取PowerPoint元数据."""
        try:
            from pptx import Presentation

            if file_path.suffix.lower() == ".pptx":
                prs = Presentation(str(file_path))

                metadata: dict[str, Any] = {
                    "slide_count": len(prs.slides),
                }

                if prs.core_properties:
                    metadata.update(
                        {
                            "title": prs.core_properties.title or "",
                            "author": prs.core_properties.author or "",
                            "subject": prs.core_properties.subject or "",
                            "created": str(prs.core_properties.created)
                            if prs.core_properties.created
                            else "",
                            "modified": str(prs.core_properties.modified)
                            if prs.core_properties.modified
                            else "",
                        }
                    )

                return metadata

            return {}

        except Exception as e:
            logger.error(f"获取PowerPoint元数据失败: {str(e)}")
            return {}

    async def split_document(
        self, content: str, chunk_size: int = 1000, overlap: int = 100
    ) -> list[str]:
        """将文档内容分割成块."""
        if not content:
            return []

        chunks = []
        start = 0

        while start < len(content):
            end = start + chunk_size

            # 如果不是最后一块，尝试在句号或换行符处分割
            if end < len(content):
                # 寻找最近的句号或换行符
                for i in range(end, max(start + chunk_size // 2, start), -1):
                    if content[i] in ".。\n":
                        end = i + 1
                        break

            chunk = content[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - overlap if end < len(content) else end

        return chunks

    async def extract_keywords(self, content: str, max_keywords: int = 10) -> list[str]:
        """提取关键词（简化版）."""
        if not content:
            return []

        # 简单的关键词提取（基于词频）
        import re
        from collections import Counter

        # 清理文本
        text = re.sub(r"[^\w\s]", " ", content.lower())
        words = text.split()

        # 过滤停用词（简化版）
        stop_words = {
            "the",
            "a",
            "an",
            "and",
            "or",
            "but",
            "in",
            "on",
            "at",
            "to",
            "for",
            "of",
            "with",
            "by",
            "的",
            "了",
            "在",
            "是",
            "我",
            "有",
            "和",
            "就",
            "不",
            "人",
            "都",
            "一",
            "一个",
            "上",
            "也",
            "很",
            "到",
            "说",
            "要",
            "去",
            "你",
            "会",
            "着",
            "没有",
            "看",
            "好",
            "自己",
            "这",
        }

        # 过滤短词和停用词
        filtered_words = [
            word for word in words if len(word) > 2 and word not in stop_words
        ]

        # 统计词频
        word_counts = Counter(filtered_words)

        # 返回最常见的词
        return [word for word, count in word_counts.most_common(max_keywords)]

    async def summarize_content(self, content: str, max_length: int = 200) -> str:
        """生成内容摘要（简化版）."""
        if not content:
            return ""

        if len(content) <= max_length:
            return content

        # 简单的摘要生成：取前几句话
        sentences = content.split("。")
        summary = ""

        for sentence in sentences:
            if len(summary + sentence + "。") <= max_length:
                summary += sentence + "。"
            else:
                break

        return summary or content[:max_length] + "..."

    async def extract_content_enhanced(
        self, file_path: Path, file_ext: str | None = None
    ) -> dict[str, Any]:
        """增强的内容提取，返回结构化数据.

        Returns:
            包含以下字段的字典:
            - text: 提取的文本内容
            - format: 内容格式 (markdown/plain_text)
            - extraction_method: 使用的提取方法
            - metadata: 额外的元数据
            - has_tables: 是否包含表格
            - has_images: 是否包含图片（仅用于标记）
            - has_formulas: 是否包含公式（仅用于标记）
        """
        if not file_ext:
            file_ext = file_path.suffix.lower()

        result: dict[str, Any] = {
            "text": "",
            "format": "plain_text",
            "extraction_method": "default",
            "metadata": {},
            "has_tables": False,
            "has_images": False,
            "has_formulas": False,
        }

        # 如果markitdown可用且支持该格式，优先使用
        if self.markitdown and file_ext in ['.pdf', '.docx', '.pptx', '.xlsx']:
            try:
                markitdown_result = self.markitdown.convert(str(file_path))
                if markitdown_result and markitdown_result.text_content:
                    result["text"] = markitdown_result.text_content
                    result["format"] = "markdown"
                    result["extraction_method"] = "markitdown"

                    # 简单检测内容特征
                    content_lower = result["text"].lower()
                    result["has_tables"] = "|" in result["text"] or "table" in content_lower
                    result["has_images"] = "![" in result["text"] or "image" in content_lower
                    result["has_formulas"] = "$$" in result["text"] or "\\begin{" in result["text"]

                    return result
            except Exception as e:
                logger.warning(f"MarkItDown extraction failed: {e}")

        # 使用默认提取器
        result["text"] = await self.extract_content(file_path, file_ext)

        # 获取元数据
        try:
            result["metadata"] = await self.get_document_metadata(file_path)
        except Exception as e:
            logger.warning(f"Failed to extract metadata: {e}")

        return result
